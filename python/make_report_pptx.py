#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
make_report_pptx.py — Builds the PowerPoint report (WHITE background) for the FPGA DV-QKD lab

Content follows `docs/bao_cao_lab_uwoc_v2.md`, numbers taken from
`uwoc_channel_model.py` / RTL, images taken from `python/Images/report/`.

Requires: pip install python-pptx
Run:      python python/make_report_pptx.py
Output:   docs/BaoCao_FPGA_DVQKD_UWOC.pptx
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG = os.path.join(ROOT, "python", "Images", "report")
OUTDIR = os.path.join(ROOT, "docs")
os.makedirs(OUTDIR, exist_ok=True)
OUTFILE = os.path.join(OUTDIR, "BaoCao_FPGA_DVQKD_UWOC.pptx")

# ── colour palette ─────────────────────────────────────────────────────────
RED = RGBColor(0xC4, 0x1E, 0x2D)        # HUST red
GOLD = RGBColor(0xE8, 0xA3, 0x3D)
INK = RGBColor(0x14, 0x14, 0x14)
INK2 = RGBColor(0x55, 0x55, 0x52)
INK3 = RGBColor(0x8A, 0x8A, 0x86)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
AQUA = RGBColor(0x16, 0x8F, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xF4, 0xF4, 0xF2)
NOTEBG = RGBColor(0xFD, 0xF3, 0xF2)

FONT = "Segoe UI"
MONO = "Consolas"

SW, SH = Inches(13.333), Inches(7.5)
ML = Inches(0.62)                        # left margin
MW = Inches(12.09)                       # content area width
BODY_TOP = Inches(1.45)

FOOTER = "FPGA-based DV-QKD  ·  Real-time UWOC BB84 emulator  ·  HUST"

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

_page = {"n": 0}


# ═══════════════════════════════════════════════════════════════════════════
# utilities
# ═══════════════════════════════════════════════════════════════════════════
def _txbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def _rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def _run(p, text, size=16, bold=False, color=INK, font=FONT, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def new_slide(title=None, subtitle=None, footer=True):
    """White-background slide: title + red rule + footer."""
    s = prs.slides.add_slide(BLANK)
    _page["n"] += 1
    if title is not None:
        _, tf = _txbox(s, ML, Inches(0.42), MW, Inches(0.62))
        p = tf.paragraphs[0]
        _run(p, title, size=27, bold=True, color=INK)
        _rect(s, ML, Inches(1.12), Inches(1.55), Inches(0.055), RED)
        if subtitle:
            _, tf2 = _txbox(s, ML, Inches(1.20), MW, Inches(0.34))
            _run(tf2.paragraphs[0], subtitle, size=13, color=INK2, italic=True)
    if footer:
        _, tf = _txbox(s, ML, Inches(7.02), Inches(9.5), Inches(0.28))
        _run(tf.paragraphs[0], FOOTER, size=9, color=INK3)
        _, tf = _txbox(s, Inches(11.9), Inches(7.02), Inches(0.81), Inches(0.28))
        pp = tf.paragraphs[0]
        pp.alignment = PP_ALIGN.RIGHT
        _run(pp, str(_page["n"]), size=10, bold=True, color=RED)
    return s


def bullets(slide, items, left=ML, top=None, width=MW, height=Inches(5.2),
            size=16, gap=8):
    """
    items: a string, or (text, level), or (text, level, dict(style)).
    level 0 = main bullet, 1 = sub-bullet, 2 = small note.
    """
    top = BODY_TOP if top is None else top
    _, tf = _txbox(slide, left, top, width, height)
    first = True
    for it in items:
        opt = {}
        if isinstance(it, str):
            text, lvl = it, 0
        elif len(it) == 2:
            text, lvl = it
        else:
            text, lvl, opt = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(opt.get("gap", gap))
        p.space_before = Pt(opt.get("gap_before", 0))
        sz = opt.get("size", size - lvl * 1.5)
        col = opt.get("color", INK if lvl == 0 else INK2)
        bold = opt.get("bold", False)
        p.level = lvl
        if text == "":
            _run(p, " ", size=6)
            continue
        marker = opt.get("marker", {0: "▪  ", 1: "–  ", 2: ""}[lvl])
        if marker:
            _run(p, marker, size=sz, bold=bold,
                 color=opt.get("marker_color", RED if lvl == 0 else INK3))
        _run(p, text, size=sz, bold=bold, color=col,
             font=opt.get("font", FONT), italic=opt.get("italic", False))
    return tf


def picture(slide, name, left, top, max_w, max_h):
    """Insert an image, keeping its aspect ratio, centred in (left, top, max_w, max_h)."""
    path = os.path.join(IMG, name)
    pic = slide.shapes.add_picture(path, left, top)
    ar = pic.width / pic.height
    w, h = max_w, int(max_w / ar)
    if h > max_h:
        h, w = max_h, int(max_h * ar)
    pic.width, pic.height = int(w), int(h)
    pic.left = int(left + (max_w - w) / 2)
    pic.top = int(top + (max_h - h) / 2)
    return pic


def fig_slide(title, name, takeaways=(), subtitle=None, side=None, note=None,
              img_h=None):
    """
    Image slide. side=None → automatic: a wide image (aspect > 2.1) goes on top with
    text below; a squarer image goes on the left with text on the right.
    """
    s = new_slide(title, subtitle)
    top = Inches(1.62) if not subtitle else Inches(1.72)
    from PIL import Image  # noqa: PLC0415
    try:
        with Image.open(os.path.join(IMG, name)) as im:
            ar = im.width / im.height
    except Exception:
        ar = 1.8
    wide = ar > 2.1 if side is None else (side == "top")

    if wide or not takeaways:
        ih = img_h or (Inches(4.05) if takeaways else Inches(5.05))
        picture(s, name, ML, top, MW, ih)
        if takeaways:
            bullets(s, takeaways, top=Inches(5.85), height=Inches(1.15),
                    size=14, gap=4)
    else:
        picture(s, name, ML, top, Inches(7.55), Inches(4.95))
        bullets(s, takeaways, left=Inches(8.45), top=top,
                width=Inches(4.25), height=Inches(4.95), size=14, gap=10)
    if note:
        _, tf = _txbox(s, ML, Inches(6.72), Inches(11.0), Inches(0.3))
        _run(tf.paragraphs[0], note, size=10.5, color=INK3, italic=True)
    return s


def table(slide, rows, left, top, width, col_w=None, size=12,
          head_size=12, row_h=Inches(0.34), head_bg=RED, zebra=True):
    """rows[0] is the header row."""
    nr, nc = len(rows), len(rows[0])
    shp = slide.shapes.add_table(nr, nc, left, top, width,
                                 Emu(int(row_h * nr))).table
    if col_w:
        tot = sum(col_w)
        for i, c in enumerate(col_w):
            shp.columns[i].width = Emu(int(width * c / tot))
    for i, row in enumerate(rows):
        shp.rows[i].height = row_h
        for j, val in enumerate(row):
            cell = shp.cell(i, j)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = 0
            cell.fill.solid()
            if i == 0:
                cell.fill.fore_color.rgb = head_bg
            else:
                cell.fill.fore_color.rgb = (
                    RGBColor(0xF7, 0xF7, 0xF5) if (zebra and i % 2 == 0)
                    else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            _run(p, str(val), size=head_size if i == 0 else size,
                 bold=(i == 0 or j == 0),
                 color=WHITE if i == 0 else INK,
                 font=FONT)
    return shp


def code_block(slide, lines, left, top, width, height, size=12.5,
               bg=RGBColor(0xF6, 0xF7, 0xF9), accent=BLUE):
    """Monospaced code block. The font size shrinks so the longest line never wraps."""
    _rect(slide, left, top, width, height, bg)
    _rect(slide, left, top, Inches(0.05), height, accent)
    inner_pt = (width - Inches(0.64)) / 914400 * 72
    longest = max((len(x) for x in lines), default=1)
    # Consolas: character width ≈ 0.55 × font size
    size = min(size, inner_pt / max(longest, 1) / 0.55)
    _, tf = _txbox(slide, left + Inches(0.24), top + Inches(0.14),
                   width - Inches(0.4), height - Inches(0.2))
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        _run(p, ln, size=size, font=MONO, color=INK)
    return tf


def callout(slide, text, left, top, width, height=Inches(0.85),
            bg=NOTEBG, accent=RED, size=13.5, label=None):
    _rect(slide, left, top, width, height, bg)
    _rect(slide, left, top, Inches(0.055), height, accent)
    _, tf = _txbox(slide, left + Inches(0.26), top + Inches(0.14),
                   width - Inches(0.45), height - Inches(0.2))
    p = tf.paragraphs[0]
    if label:
        _run(p, label + "  ", size=size, bold=True, color=accent)
    _run(p, text, size=size, color=INK)
    return tf


def section(number, title, points=()):
    s = new_slide(None, footer=True)
    _rect(s, Inches(0), Inches(2.35), SW, Inches(2.35), BAND)
    _rect(s, Inches(0), Inches(2.35), Inches(0.16), Inches(2.35), RED)
    _, tf = _txbox(s, Inches(0.95), Inches(2.62), Inches(11.5), Inches(0.5))
    _run(tf.paragraphs[0], f"PART {number}", size=15, bold=True, color=RED)
    _, tf = _txbox(s, Inches(0.95), Inches(3.10), Inches(11.5), Inches(0.8))
    _run(tf.paragraphs[0], title, size=34, bold=True, color=INK)
    if points:
        _, tf = _txbox(s, Inches(0.95), Inches(4.02), Inches(11.5), Inches(0.5))
        _run(tf.paragraphs[0], "   ·   ".join(points), size=13.5, color=INK2)
    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(None, footer=False)
_page["n"] = 0
_rect(s, Inches(0), Inches(0), SW, Inches(0.22), RED)
_rect(s, Inches(0), Inches(0.22), SW, Inches(0.06), GOLD)

_, tf = _txbox(s, ML, Inches(2.05), Inches(11.6), Inches(1.35))
_run(tf.paragraphs[0], "FPGA-BASED DV-QKD", size=52, bold=True, color=RED)

_, tf = _txbox(s, ML, Inches(3.30), Inches(11.6), Inches(1.0))
p = tf.paragraphs[0]
_run(p, "Real-Time Hardware Emulation of a BB84-QKD System\n", size=21,
     color=INK)
p2 = tf.add_paragraph()
_run(p2, "in an Underwater Optical Channel, with Closed-Loop Adaptive Control",
     size=21, color=INK)

_rect(s, ML, Inches(4.55), Inches(2.1), Inches(0.045), GOLD)

_, tf = _txbox(s, ML, Inches(4.85), Inches(6.0), Inches(1.0))
_run(tf.paragraphs[0], "Lê Công Khánh", size=17, color=INK)
p = tf.add_paragraph()
_run(p, "Đào Đặng Minh Hoàng", size=17, color=INK)

_, tf = _txbox(s, ML, Inches(6.15), Inches(8.0), Inches(0.7))
_run(tf.paragraphs[0], "Hanoi University of Science and Technology",
     size=13, color=INK2)
p = tf.add_paragraph()
_run(p, "School of Electrical and Electronic Engineering", size=13, color=INK2)

_, tf = _txbox(s, Inches(9.0), Inches(6.15), Inches(3.7), Inches(0.7))
pp = tf.paragraphs[0]
pp.alignment = PP_ALIGN.RIGHT
_run(pp, "Lab report — simulation results", size=13, bold=True, color=RED)
p = tf.add_paragraph()
p.alignment = PP_ALIGN.RIGHT
_run(p, "August 2026", size=13, color=INK2)

_rect(s, Inches(0), Inches(7.28), SW, Inches(0.22), RED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Contents
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Outline")
items = [
    ("I", "Objectives, scope and motivation",
     "what the emulator is — and is not"),
    ("II", "The underwater optical channel",
     "path loss, turbulence, scattering, depolarization"),
    ("III", "FPGA implementation",
     "ROM tables, detection layer, monitor, adaptive controller"),
    ("IV", "Simulation results",
     "QBER, P_click, SKR, wavelength, turbulence"),
    ("V", "Verification, status and future work",
     "what is verified, what is pending"),
]
_, tf = _txbox(s, ML, Inches(1.70), MW, Inches(4.9))
for i, (num, title, sub) in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(16)
    _run(p, f"{num:>3}   ", size=19, bold=True, color=RED, font=MONO)
    _run(p, title + "     ", size=19, bold=True, color=INK)
    _run(p, sub, size=14, color=INK2, italic=True)

callout(s, "The board fault is not yet resolved, so every number in Part IV "
           "is a Monte-Carlo simulation of the channel model that is loaded "
           "into the FPGA ROM — not a hardware measurement.",
        ML, Inches(6.15), MW, Inches(0.82), label="Status:")

# ═══════════════════════════════════════════════════════════════════════════
section("I", "Objectives, Scope and Motivation",
        ("Topic", "Objectives", "What this work does not claim"))

# ── Topic & objectives ────────────────────────────────────────────────────
s = new_slide("Topic and Objectives")
bullets(s, [
    ("Topic", 0, dict(bold=True, size=17)),
    ("Real-Time Hardware Emulation of a BB84-QKD System in an Underwater "
     "Optical Channel, with Closed-Loop Adaptive Control", 1, dict(size=15)),
    ("", 0),
    ("Objectives", 0, dict(bold=True, size=17)),
    ("Emulate the UWOC channel in real time on a single FPGA: deterministic "
     "path loss, scattering fading, oceanic turbulence, distance-dependent "
     "polarization error, detector noise", 1, dict(size=14.5)),
    ("Implement BB84 transmission, sifting and per-qubit error detection",
     1, dict(size=14.5)),
    ("Estimate QBER, SNR and turbulence jitter on-chip over a 65,536-attempt "
     "window, and close an adaptive control loop on those estimates",
     1, dict(size=14.5)),
    ("Adapt four transmission parameters in real time: optical drive level, "
     "basis bias, slot width, and wavelength", 1, dict(size=14.5)),
    ("Stream per-qubit results to a PC over UART for statistical analysis",
     1, dict(size=14.5)),
], height=Inches(4.0))

callout(s, "This is an FPGA-based EMULATOR, not a physical QKD terminal. "
           "There is no laser, no single-photon detector and no water tank. "
           "The FPGA reproduces the channel statistics and the BB84 reception "
           "process in real time. QBER and SKR reported here are computed on "
           "the PC; the on-chip estimator serves the control loop.",
        ML, Inches(5.60), MW, Inches(1.15), label="Scope note:")

# ── Motivation ────────────────────────────────────────────────────────────
s = new_slide("Motivation: why underwater optical QKD?")
bullets(s, [
    ("Underwater sensor networks, AUVs and marine monitoring systems need "
     "secure links. QKD gives key distribution whose security rests on "
     "physics, not on computational hardness.", 2, dict(size=15)),
    ("", 0),
    ("Acoustic communication", 0, dict(bold=True, size=16)),
    ("✓  range of kilometres to tens of kilometres", 1, dict(size=14)),
    ("✗  bandwidth only in the kbps range;  ✗  latency — sound travels at "
     "1500 m/s", 1, dict(size=14)),
    ("✗  cannot carry photon polarization states, so BB84 is impossible",
     1, dict(size=14, color=RED)),
    ("", 0),
    ("Underwater optical communication (UWOC)", 0, dict(bold=True, size=16)),
    ("✓  Mbps to Gbps;  ✓  low latency;  ✓  photons carry the BB84 quantum "
     "states directly", 1, dict(size=14)),
    ("✗  severe attenuation limits the range to metres or tens of metres",
     1, dict(size=14)),
], height=Inches(4.4))

callout(s, "BB84 is already fragile against loss and noise. Underwater the "
           "dominant event is not a bit flip — it is a MISSING PHOTON: "
           "P_click runs from 1e-2 down to 1e-5 over the distances of "
           "interest. Both the hardware model and the measurement "
           "methodology have to be built around that fact.",
        ML, Inches(5.72), MW, Inches(1.05), label="Consequence:")

# ── Problem statement / scope ─────────────────────────────────────────────
s = new_slide("Problem statement and system scope")
bullets(s, [("Model inputs", 0, dict(bold=True, size=16))],
        top=Inches(1.50), height=Inches(0.4))
table(s, [
    ["Parameter", "Values"],
    ["Water type", "Clear ocean  ·  Coastal  ·  Turbid harbor"],
    ["Distance", "16-point grid per water type (4-bit index)"],
    ["", "clear 5–80 m  ·  coastal 2–43 m  ·  harbor 0.5–10 m"],
    ["Turbulence", "5 levels L1–L5,  σ²_ho = 0.02 / 0.08 / 0.30 / 1.00 / 3.00"],
    ["Wavelength", "450 / 532 / 650 nm"],
    ["Source, detector", "μ = 0.1 photon/pulse, η_det = 0.18, Y₀ = 1.3e-5, "
                         "f_rep = 10 MHz"],
], ML, Inches(1.95), MW, col_w=[26, 74], size=13, head_size=13,
    row_h=Inches(0.36))

bullets(s, [("Outputs", 0, dict(bold=True, size=16))],
        top=Inches(4.55), height=Inches(0.4))
bullets(s, [
    ("Measurement mode — one record per qubit over UART:", 1, dict(size=14)),
    ("@a_data, a_basis, b_basis, bob_bit, bmatch, err, irrad, total, "
     "sifted, errors*", 2, dict(size=12.5, font=MONO, color=BLUE)),
    ("Control mode — on-chip window estimates: QBER, SNR proxy, QBER jitter, "
     "loss rate, window_valid", 1, dict(size=14)),
], top=Inches(4.98), height=Inches(1.5), gap=5)

# ── Block diagram ─────────────────────────────────────────────────────────
fig_slide("System architecture", "fig00_system_block.png", side="top",
          subtitle="all BB84 generation, channel emulation, detection, sifting "
                   "and monitoring run on ONE FPGA — the PC only logs and plots",
          note="Clock 50 MHz  ·  UART 115,200 bps  ·  fading resampled every "
               "2¹⁸ cycles  ·  h: 12-bit (256 = 1.0)  ·  probabilities: 24-bit "
               "·  monitoring window 2¹⁶ attempts  ·  MIN_SIFT = 16")

# ═══════════════════════════════════════════════════════════════════════════
section("II", "The Underwater Optical Channel",
        ("Physics", "Composite model", "Path loss", "Turbulence", "Scattering",
         "Depolarization"))

s = fig_slide("Physical impairments of the underwater channel",
              "fig00_channel_physics.png", side="top", img_h=Inches(3.90))
callout(s, "All five effects enter the hardware model. Bubbles are currently "
           "absorbed into the scattering fading term h_s rather than modelled "
           "as a distinct blockage process — see Future work.",
        ML, Inches(5.70), MW, Inches(0.85), label="Note:", size=13.5)

# ── Combined model ────────────────────────────────────────────────────────
s = new_slide("Composite channel model")
code_block(s, [
    "h(d, t)  =  L(d, λ, water)  ·  h_s(d, t)  ·  h_o(d, t)",
    "",
    "     L     deterministic path loss   →  Beer-Lambert + geometric spreading",
    "     h_s   scattering fading         →  Gamma,          E[h_s] = 1",
    "     h_o   oceanic turbulence        →  Log-normal (L1-L3) / Weibull (L4-L5),",
    "                                        E[h_o] = 1",
], ML, Inches(1.60), MW, Inches(1.95), size=14)

code_block(s, [
    "n̄        =  μ · L(d) · h · η_det              mean photon number at the detector",
    "P_click  =  1 − (1 − Y₀) · exp(−n̄)",
], ML, Inches(3.72), MW, Inches(0.95), size=14, accent=AQUA,
    bg=RGBColor(0xF2, 0xFA, 0xF6))

bullets(s, [
    ("Both fading terms are normalised to unit mean, so all deterministic "
     "loss lives in L and the two random terms only redistribute it in time.",
     0, dict(size=14.5)),
    ("Unlike the previous report, μ and η_det are NOT folded into an "
     "\"L_eff\": μ is a knob the adaptive controller moves (with an anti-PNS "
     "cap), while η_det is a fixed detector constant. Merging them would hide "
     "the control loop.", 0, dict(size=14.5)),
], top=Inches(4.95), height=Inches(1.6), gap=9)

_, tf = _txbox(s, ML, Inches(6.60), MW, Inches(0.3))
_run(tf.paragraphs[0],
     "Reference: Salcedo-Serrano et al., IEEE ICC 2022, eq. (2).",
     size=11, color=INK3, italic=True)

# ── Path loss ─────────────────────────────────────────────────────────────
s = new_slide("Deterministic path loss")
code_block(s, [
    "                D_rx²",
    "  L(d)  =  ───────────────────  ·  exp[ −F · c(λ) · d ]",
    "            π (d · tan θ_div)²",
    "            └─ geometric ─┘        └── Beer-Lambert ──┘",
], ML, Inches(1.62), Inches(6.9), Inches(1.75), size=13)

bullets(s, [
    ("D_rx = 50.8 mm  —  receiver aperture (2 inch)", 1, dict(size=13.5)),
    ("θ_div = 1e-3 rad  —  effective half-divergence; diffraction alone is "
     "~4e-5 rad, the rest is pointing-error margin", 1, dict(size=13.5)),
    ("F = 0.85  —  fraction of scattered photons still collected",
     1, dict(size=13.5)),
    ("c(λ) = a(λ) + b(λ)  —  extinction coefficient", 1, dict(size=13.5)),
    ("The geometric term is clamped to ≤ 1: the receiver cannot collect more "
     "than was transmitted.", 1, dict(size=13.5, color=INK2, italic=True)),
], top=Inches(3.58), width=Inches(6.25), height=Inches(1.9), gap=6)

picture(s, "fig01_pathloss.png", Inches(7.05), Inches(1.62),
        Inches(5.68), Inches(4.55))
_, tf = _txbox(s, ML, Inches(5.72), Inches(6.6), Inches(1.0))
_run(tf.paragraphs[0],
     "Over the deployed grids the three water types span eleven orders of "
     "magnitude in loss. This is why each water type needs its own distance "
     "grid rather than a single 1–7 m range.", size=13, color=INK)

# ── Water types & link budget ─────────────────────────────────────────────
s = new_slide("Water types and link budget")
table(s, [
    ["Water type", "a [1/m]", "b [1/m]", "c [1/m]", "Distance grid"],
    ["Clear ocean", "0.114", "0.037", "0.151", "5 – 80 m  (step 5 m)"],
    ["Coastal", "0.179", "0.219", "0.398", "2 – 43 m"],
    ["Turbid harbor", "0.366", "1.824", "2.190", "0.5 – 10 m"],
], ML, Inches(1.52), Inches(6.55), col_w=[26, 17, 17, 17, 30], size=12.5,
    head_size=12, row_h=Inches(0.35))

_, tf = _txbox(s, ML, Inches(3.10), Inches(6.55), Inches(0.3))
_run(tf.paragraphs[0], "Wavelength scaling of c(λ), relative to 532 nm",
     size=13, bold=True, color=INK)
table(s, [
    ["Water type", "450 nm", "532 nm", "650 nm"],
    ["Clear ocean", "0.85", "1.00", "2.60"],
    ["Coastal", "1.10", "1.00", "1.45"],
    ["Turbid harbor", "1.35", "1.00", "0.90"],
], ML, Inches(3.48), Inches(6.55), col_w=[34, 22, 22, 22], size=12.5,
    head_size=12, row_h=Inches(0.35))

bullets(s, [("Link budget", 0, dict(bold=True, size=15))],
        left=Inches(7.35), top=Inches(1.52), width=Inches(5.4),
        height=Inches(0.4))
code_block(s, [
    "μ      = 0.1 photon/pulse   weak coherent source",
    "η_det  = 0.18               single-photon detector",
    "Y₀     = (60 Hz + 200 Hz) × 50 ns = 1.3e-5",
    "f_rep  = 10 MHz",
    "τ_coh  = 5 ms  →  50,000 pulses per block",
], Inches(7.35), Inches(1.96), Inches(5.4), Inches(1.62), size=12)

bullets(s, [("Fixed-point formats on the FPGA", 0, dict(bold=True, size=15))],
        left=Inches(7.35), top=Inches(3.78), width=Inches(5.4),
        height=Inches(0.4))
code_block(s, [
    "h              12-bit, 256 = 1.0, saturates at 4095",
    "probabilities  24-bit, value = round(p × 2²⁴)",
], Inches(7.35), Inches(4.22), Inches(5.4), Inches(0.86), size=12,
    accent=AQUA, bg=RGBColor(0xF2, 0xFA, 0xF6))

callout(s, "The λ-scaling table is INDICATIVE, not a measured spectrum. The "
           "qualitative ordering follows Rosenkrantz & Arnon; a precise "
           "crossover would need measured a(λ), b(λ) per water type.",
        ML, Inches(5.55), MW, Inches(0.80), label="Caveat:", size=13)

# ── Turbulence L1–L5 ──────────────────────────────────────────────────────
s = new_slide("Oceanic turbulence: five calibrated levels")
bullets(s, [
    ("Turbulence strength is not a free knob. Each level is a physically "
     "consistent triple (ε, χ_T, w) pushed through the Nikishov spatial power "
     "spectrum, then calibrated so the resulting scintillation index hits a "
     "target at d = 20 m, λ = 450 nm.", 2, dict(size=14)),
], top=Inches(1.45), height=Inches(0.85))

table(s, [
    ["Level", "ε [m²/s³]", "χ_T [K²/s]", "w", "σ²_ho (20 m)", "Distribution"],
    ["L1  Very weak", "1e-2", "2.21e-7", "−5.0", "0.02", "Log-normal"],
    ["L2  Weak", "1e-3", "3.65e-7", "−4.0", "0.08", "Log-normal"],
    ["L3  Moderate", "1e-4", "5.29e-7", "−3.0", "0.30", "Log-normal"],
    ["L4  Strong", "1e-5", "5.91e-7", "−2.0", "1.00", "Weibull"],
    ["L5  Severe", "1e-6", "3.85e-7", "−1.0", "3.00", "Weibull"],
], ML, Inches(2.32), Inches(6.80), col_w=[25, 15, 16, 11, 17, 20],
    size=11.5, head_size=11, row_h=Inches(0.345))

bullets(s, [
    ("ε — turbulent kinetic energy dissipation rate", 2, dict(size=13)),
    ("χ_T — dissipation rate of mean-square temperature", 2, dict(size=13)),
    ("w — temperature/salinity contribution ratio", 2, dict(size=13)),
    ("      w = −5 salinity-driven … w = −1 temperature-driven",
     2, dict(size=12, color=INK3)),
    ("", 0),
    ("L4 = 1.0 is the conventional weak/strong turbulence boundary.",
     2, dict(size=13, bold=True)),
    ("σ²_ho is clamped at 5.0 to keep the ROM inversion numerically stable — "
     "the Rytov result is only valid for σ² ≲ 1.", 2, dict(size=13)),
], top=Inches(4.55), width=Inches(6.80), height=Inches(1.9), gap=4)

picture(s, "fig02_sigma2_ho.png", Inches(7.60), Inches(2.05),
        Inches(5.15), Inches(4.35))

# ── Log-normal / Weibull ──────────────────────────────────────────────────
s = new_slide("Turbulence fading distribution",
              "the distribution is chosen by turbulence strength, not by water type")
code_block(s, [
    "Weak to moderate turbulence  (L1–L3)   →   Log-normal",
    "",
    "    f(h_o) = 1/(h_o σ √2π) · exp( −(ln h_o − μ)² / 2σ² )",
    "    μ = −σ²/2      so that   E[h_o] = 1",
], ML, Inches(1.85), Inches(6.0), Inches(1.75), size=12.5)

code_block(s, [
    "Strong turbulence  (L4–L5)             →   Weibull",
    "",
    "    f(h_o) = (β/λ)(h_o/λ)^(β−1) · exp( −(h_o/λ)^β )",
    "    λ = 1 / Γ(1 + 1/β)   so that   E[h_o] = 1",
    "    smaller β  =  deeper fades",
], Inches(6.85), Inches(1.85), Inches(5.9), Inches(1.75), size=12.5,
    accent=RGBColor(0x17, 0x49, 0x7F), bg=RGBColor(0xF2, 0xF6, 0xFB))

bullets(s, [
    ("Both are stored as 256-point inverse-CDF tables per level "
     "(8 × 256 × 12 bit) — sampling a distribution on the FPGA is one ROM "
     "read indexed by an LFSR word.", 0, dict(size=15)),
], top=Inches(3.95), height=Inches(0.8))

callout(s, "The previous report tied Weibull to \"turbid harbor\" and "
           "log-normal to \"coastal water\". That is a conceptual error: water "
           "type determines ATTENUATION and h_s; turbulence is an independent "
           "variable, and the same body of water can sit at any level L1–L5.",
        ML, Inches(4.85), MW, Inches(1.15), label="Correction:")

# ── Scattering fading ─────────────────────────────────────────────────────
s = new_slide("Scattering-induced fading  h_s")
bullets(s, [
    ("Multiple scattering in water spreads the beam in time and angle, so the "
     "collected power fluctuates. Modelled as a Gamma variable with unit mean:",
     2, dict(size=15)),
], top=Inches(1.48), height=Inches(0.6))

code_block(s, [
    "  f(h_s; k, θ)  =  h_s^(k−1) · e^(−h_s/θ)  /  ( Γ(k) · θ^k )",
    "",
    "  k = 1/σ_s²       θ = σ_s²        ⇒   E[h_s] = k·θ = 1",
], ML, Inches(2.15), Inches(7.6), Inches(1.45), size=13)

bullets(s, [
    ("σ_s² grows with distance and turbidity; it is fitted from Table IV of "
     "Salcedo-Serrano ICC'22 as  ln σ_s² = B (d − d₁).", 0, dict(size=14.5)),
    ("Hardware: 8 fading classes  (σ_s² = 0, 0.02, 0.05, 0.12, 0.3, 0.75, "
     "1.8, 4.5). Each (water, distance) cell selects a class through "
     "hscls_rom.", 0, dict(size=14.5)),
], top=Inches(3.85), height=Inches(1.5), gap=9)

callout(s, "Clear ocean and coastal are fitted to published measurements. "
           "Harbor is EXTRAPOLATED by a power law in c(λ) — the source paper "
           "publishes no harbor data. The report states this rather than "
           "presenting all three rows as equally grounded.",
        ML, Inches(5.35), MW, Inches(1.05), label="Caveat:")

# ── e_pol ─────────────────────────────────────────────────────────────────
s = new_slide("Polarization error — the UWOC-specific QBER mechanism")
bullets(s, [
    ("In free-space optical links QBER is dominated by deep fades and "
     "detector noise. Underwater there is an extra mechanism with no "
     "atmospheric counterpart: multiple scattering DEPOLARIZES the photon, so "
     "a photon that arrives can still land in the wrong polarization gate.",
     2, dict(size=14)),
], top=Inches(1.45), height=Inches(0.9))

code_block(s, [
    "e_pol(d) = e₀ + k_s · [ 1 − exp(−b(λ)·d) ]    capped at 0.5",
    "",
    "   e₀  = 0.01   intrinsic optical error",
    "                (misalignment, PBS extinction)",
    "   k_s = 0.04   scattering depolarization coefficient",
    "   b            scattering coefficient at λ",
], ML, Inches(2.30), Inches(7.15), Inches(2.10), size=12.5)

code_block(s, [
    "           e_pol · (1 − e^(−n̄))  +  0.5 · Y₀",
    "QBER  =  ───────────────────────────────────",
    "              1 − (1 − Y₀) · e^(−n̄)",
], ML, Inches(4.52), Inches(7.15), Inches(1.30), size=12.5, accent=AQUA,
    bg=RGBColor(0xF2, 0xFA, 0xF6))

bullets(s, [
    ("Signal clicks are wrong with probability e_pol; dark and background "
     "clicks carry no information, so they are wrong half the time.",
     2, dict(size=12.5)),
    ("For small b·d this is linear in d, matching the linear QBER-vs-distance "
     "trend reported by Kebapci et al. (2023, Fig. 19).",
     2, dict(size=12.5)),
], top=Inches(5.94), width=Inches(7.15), height=Inches(1.0), gap=4)

picture(s, "fig03_epol.png", Inches(7.90), Inches(2.30),
        Inches(4.85), Inches(3.35))
callout(s, "k_s is a CALIBRATION parameter and has not yet been fitted "
           "against measured data — a known weakness, stated deliberately.",
        Inches(7.90), Inches(5.80), Inches(4.85), Inches(1.05),
        label="Caveat:", size=12.5)

# ═══════════════════════════════════════════════════════════════════════════
section("III", "FPGA Implementation",
        ("Detection layer", "Coherence-time sampling", "ROM architecture",
         "Monitor", "Adaptive control"))

# ── Decision stage ────────────────────────────────────────────────────────
s = new_slide("Detection decision layer — division-free by construction")
bullets(s, [
    ("A 50 MHz Cyclone II cannot divide once per photon event. The QBER "
     "expression is a RATIO — but it never has to be evaluated as one. "
     "Sampling numerator and denominator separately is exactly equivalent:",
     2, dict(size=14.5)),
], top=Inches(1.45), height=Inches(0.85))

code_block(s, [
    "  sig_det    =  rand24 < p_sig                  signal-induced click",
    "  noise_det  =  rand24 < p_noise                dark count + background",
    "  click      =  sig_det  OR  noise_det",
    "  err        =  sig_det ? (rand24 < e_pol) : rand24[23]",
], ML, Inches(2.32), MW, Inches(1.75), size=14)

bullets(s, [
    ("The last line carries the physics: a click caused by the signal is "
     "wrong with probability e_pol; a click caused only by background carries "
     "no information at all, so it is wrong with probability ½ — one random "
     "bit.", 0, dict(size=14.5)),
    ("Averaged over many qubits this reproduces the closed-form QBER exactly, "
     "using four comparators and no divider.", 0, dict(size=14.5)),
    ("p_noise = 218 / 2²⁴ = 1.30e-5, matching Y₀ of the analytical model.",
     0, dict(size=14.5)),
    ("Five independent 32-bit LFSRs drive h_s, h_o, signal, noise and error. "
     "Channel emulator only — BB84 entropy comes from the ring-oscillator "
     "TRNG.", 0, dict(size=14.5)),
], top=Inches(4.35), height=Inches(2.4), gap=9)

# ── Coherence sampling ────────────────────────────────────────────────────
fig_slide("Real-time fading: coherence-time block sampling",
          "fig11_coherence_timing.png",
          takeaways=[
              ("Oceanic turbulence evolves on a millisecond scale — about 10⁴ "
               "times slower than the qubit period. h_s and h_o are re-read "
               "from ROM once every 2^COH_LOG2 = 2¹⁸ cycles = 5.24 ms ≈ τ_coh, "
               "and frozen in between.", 2, dict(size=13.5)),
              ("Resampling per qubit — as the previous atmospheric emulator "
               "did — averages the fading out inside every monitoring window "
               "and leaves the controller blind to turbulence. An IIR smoother "
               "would fix the smoothness but its time constant comes from a "
               "shift parameter, not from physics.", 2, dict(size=13, color=INK2)),
          ], side="top")

# ── ROM ───────────────────────────────────────────────────────────────────
s = new_slide("ROM-based channel: tables and addressing")
bullets(s, [
    ("All distribution sampling and all transcendental functions are "
     "precomputed offline in Python and shipped as ROM "
     "(uwoc_channel_rom.vh, generated).", 2, dict(size=14)),
], top=Inches(1.45), height=Inches(0.55))

table(s, [
    ["ROM", "Size", "Contents"],
    ["hs_rom", "8 × 256 × 12 b", "inverse CDF of h_s, per fading class"],
    ["ho_rom", "8 × 256 × 12 b", "inverse CDF of h_o, per turbulence level"],
    ["psig_rom", "144 × 24 b", "signal click probability at h = 1"],
    ["epol_rom", "144 × 24 b", "polarization error e_pol"],
    ["nexp_inv_rom", "144 × 16 b", "SNR normalization constant"],
    ["hscls_rom", "48 × 3 b", "fading class per (water, distance)"],
    ["hooff_rom", "48 × 3 b", "turbulence level offset per distance"],
], ML, Inches(2.05), Inches(7.5), col_w=[27, 25, 48], size=12,
    head_size=12, row_h=Inches(0.335))

code_block(s, [
    "wd_addr  = {water[1:0], dist[3:0]}",
    "                              48 entries",
    "lwd_addr = (λ_idx·3 + water)·16 + dist",
    "                             144 entries",
    "",
    "×3 as add-and-shift — no multiplier",
], Inches(7.85), Inches(2.05), Inches(4.9), Inches(1.8), size=11.5)

callout(s, "ROM reads MUST be registered. A combinational read makes Quartus "
           "flatten the tables into ~9.2 kbit of combinational LUTs, and the "
           "initial blocks risk being dropped — which would silently zero "
           "psig and produce an emulator that never clicks and never reports "
           "an error.",
        Inches(7.85), Inches(4.05), Inches(4.9), Inches(1.65),
        label="Lesson learned:", size=12.5)

bullets(s, [
    ("Offline Python defines the statistics; the ROM stores the samples; the "
     "coherence-block counter turns them into a smooth real-time channel.",
     0, dict(size=14, italic=True)),
], top=Inches(4.95), width=Inches(7.1), height=Inches(0.8))

# ── FPGA blocks ───────────────────────────────────────────────────────────
s = new_slide("FPGA blocks")
bullets(s, [
    ("1 · Random source — trng.v ×3", 0, dict(bold=True, size=15)),
    ("Four ring oscillators (3/5/7/9 inverters, keep attribute) XORed and "
     "Von Neumann debiased → Alice bit, Alice basis, Bob basis", 1, dict(size=13)),
    ("2 · Modulation — ook_tx_serializer.v / ook_rx_deserializer.v",
     0, dict(bold=True, size=15)),
    ("4-slot OOK frame [SYNC][basis][data][IDLE]; RX samples slot centres at "
     "1.5× and 2.5× the slot width after the detected rising edge",
     1, dict(size=13)),
    ("3 · Channel emulator — uwoc_channel.v", 0, dict(bold=True, size=15)),
    ("ROM lookup → h_s × h_o → p_sig → click / error. A lost photon silences "
     "the whole frame; an error inverts ONLY the data slot, leaving SYNC and "
     "basis intact", 1, dict(size=13)),
    ("4 · BB84 processing — bob.v, error_estimation.v", 0, dict(bold=True, size=15)),
    ("basis_match = ~(b XOR b'),   error = basis_match AND (a XOR a')",
     1, dict(size=13, font=MONO)),
], top=Inches(1.45), width=Inches(6.35), height=Inches(4.6), gap=6)

bullets(s, [
    ("5 · Channel monitor — channel_monitor.v", 0, dict(bold=True, size=15)),
    ("65,536-attempt window → QBER, SNR proxy, QBER jitter, loss rate, "
     "window_valid. Feeds the adaptive controller", 1, dict(size=13)),
    ("6 · Adaptive controller — adaptive_controller.v", 0, dict(bold=True, size=15)),
    ("Four operating modes with asymmetric hysteresis, plus wavelength "
     "hill-climbing", 1, dict(size=13)),
    ("7 · Reporting — uart_reporter.v", 0, dict(bold=True, size=15)),
    ("Measurement mode: one 34-byte record PER QUBIT. Control mode: periodic "
     "status packet", 1, dict(size=13)),
], left=Inches(7.05), top=Inches(1.45), width=Inches(5.7),
    height=Inches(3.4), gap=6)

callout(s, "Unlike the previous report, there is no on-chip entropy LUT for "
           "SKR. The controller exposes a key_rate signal used for control "
           "only; the published SKR is computed on the PC from sifted counts "
           "and their confidence interval.",
        Inches(7.05), Inches(5.05), Inches(5.7), Inches(1.3),
        label="Changed:", size=12.5)

# ── Channel monitor ───────────────────────────────────────────────────────
s = new_slide("On-chip channel monitor")
table(s, [
    ["Output", "Range", "Meaning"],
    ["qber", "0–200", "1 unit = 0.5 %, shift-compare ladder, 12 levels"],
    ["snr_level", "0–255", "128 = nominal click rate for this configuration"],
    ["loss_rate", "0–255", "255 = everything lost"],
    ["qber_jitter", "0–255", "EWMA (1/4) of |QBER(k) − QBER(k−1)|"],
    ["window_valid", "gate", "asserted only when n_sift ≥ MIN_SIFT (= 16)"],
], ML, Inches(1.55), MW, col_w=[20, 13, 67], size=13, head_size=12.5,
    row_h=Inches(0.36))

bullets(s, [
    ("Window = 2¹⁶ attempts, not 256.", 0, dict(bold=True, size=15.5)),
    ("At P_click ~ 1e-3 a 256-attempt window collects 0.4 clicks. The "
     "per-window QBER is then pure shot noise, not channel information.",
     1, dict(size=14)),
    ("snr_level is NORMALIZED by the expected click count.",
     0, dict(bold=True, size=15.5)),
    ("A raw click count underwater is dominated by exp(−c·d) — a constant of "
     "the link configuration, not something the controller can improve. "
     "Dividing it out (via nexp_inv_rom, as a multiply) makes 128 mean "
     "\"performing as expected\" regardless of distance and water type.",
     1, dict(size=14)),
    ("All five outputs are 8-bit and division-free.", 0, dict(size=14.5)),
], top=Inches(4.05), height=Inches(2.6), gap=8)

# ── 3 findings ────────────────────────────────────────────────────────────
s = new_slide("Three findings that drove the design")
for i, (head, body) in enumerate([
    ("1 · A large monitoring window is mandatory",
     "With P_click ~ 1e-3, small windows make the per-window QBER pure shot "
     "noise. At a 256-attempt window the QBER standard deviation under WEAK "
     "turbulence came out LARGER than under SEVERE turbulence — the "
     "controller would have reacted backwards."),
    ("2 · Turbulence does not show up in mean QBER",
     "Across L1–L5 the mean QBER moves a few percent, while the "
     "between-window standard deviation moves by a factor of six. P_click is "
     "nearly linear in h, so the expectation cancels the fading. Underwater "
     "turbulence is visible in the VARIANCE → qber_jitter was added; without "
     "it the controller is blind to turbulence."),
    ("3 · loss_rate cannot detect a dead link",
     "It saturates at 255 even on a perfectly healthy underwater link, simply "
     "because attenuation is high. Link-death detection uses a zero photon "
     "count over 8 consecutive windows instead."),
]):
    y = Inches(1.55 + i * 1.72)
    _rect(s, ML, y, MW, Inches(1.48), RGBColor(0xF7, 0xF8, 0xFA))
    _rect(s, ML, y, Inches(0.055), Inches(1.48), BLUE)
    _, tf = _txbox(s, ML + Inches(0.28), y + Inches(0.16),
                   MW - Inches(0.55), Inches(1.2))
    _run(tf.paragraphs[0], head, size=16, bold=True, color=BLUE)
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    _run(p, body, size=13.5, color=INK)

_, tf = _txbox(s, ML, Inches(6.72), MW, Inches(0.3))
_run(tf.paragraphs[0],
     "All three are reproducible with:  python python/uwoc_channel_model.py  "
     "(7 self-tests)", size=11.5, color=INK3, italic=True)

# ── Adaptive controller ───────────────────────────────────────────────────
s = new_slide("Adaptive controller — four operating modes")
bullets(s, [
    ("Continuous optimization is impractical on a Cyclone II. The channel "
     "state space is partitioned into four regions, each with a precomputed "
     "parameter set.", 2, dict(size=14)),
], top=Inches(1.45), height=Inches(0.55))

table(s, [
    ["Mode", "Drive", "Basis bias", "Slot", "Gap", "Entry condition"],
    ["AGGRESSIVE", "6", "128 (50 % Z)", "5 ms", "1",
     "QBER < 4 %, SNR ≥ 160, jitter < 6"],
    ["MODERATE", "9", "154 (60 % Z)", "10 ms", "2", "default"],
    ["CONSERVATIVE", "12", "204 (80 % Z)", "50 ms", "4",
     "QBER ≥ 8 % or SNR < 96 or jitter ≥ 16"],
    ["PAUSE", "—", "—", "—", "16", "QBER ≥ 15 % or SNR ≤ 40"],
], ML, Inches(2.08), MW, col_w=[17, 9, 15, 10, 8, 41], size=12.5,
    head_size=12, row_h=Inches(0.38))

bullets(s, [
    ("Jitter alone can force a downgrade even when the mean QBER still looks "
     "healthy — the direct consequence of Finding 2.", 0, dict(size=14.5)),
    ("The Z-basis bias raises sifting efficiency from 0.50 to 0.68 (+35 %), "
     "which matters more underwater than in air because every sifted bit is "
     "expensive.", 0, dict(size=14.5)),
], top=Inches(4.30), height=Inches(1.3), gap=9)

callout(s, "Downgrade (toward PAUSE) is IMMEDIATE; upgrade (toward "
           "AGGRESSIVE) requires 3 consecutive favourable windows. Reacting "
           "late to a degrading channel leaks key; reacting late to an "
           "improving channel only costs throughput.",
        ML, Inches(5.62), MW, Inches(0.95), label="Asymmetric hysteresis:")

# ── λ hill-climbing + PNS ─────────────────────────────────────────────────
s = new_slide("Wavelength hill-climbing and the anti-PNS drive cap")
code_block(s, [
    "SETTLE  →  BASE  →  WARM  →  MEAS  →  decide",
    "",
    "  BASE   accumulate click count at the current best λ",
    "  WARM   switch λ, DISCARD one window (ROM address and the",
    "         coherence block both need to settle)",
    "  MEAS   accumulate at the candidate; keep it only if it beats",
    "         the base",
], ML, Inches(1.52), Inches(6.35), Inches(2.35), size=12)

bullets(s, [
    ("Uses the 16-bit photon_count, not the 8-bit photon_rate — the 8-bit "
     "version saturates and the comparison becomes meaningless. Found by "
     "tb_adaptive_loop.v, not by inspection.", 0, dict(size=13.5)),
    ("Because the optimal λ depends on turbidity, a link that drifts into a "
     "turbid plume can recover by changing colour rather than by giving up "
     "key rate.", 0, dict(size=13.5)),
], top=Inches(4.05), width=Inches(6.35), height=Inches(1.9), gap=8)

_, tf = _txbox(s, Inches(7.05), Inches(1.52), Inches(5.7), Inches(0.4))
_run(tf.paragraphs[0], "Why the drive level is capped", size=16, bold=True,
     color=INK)
bullets(s, [
    ("With a weak coherent source, the fraction of non-empty pulses carrying "
     "two or more photons is", 2, dict(size=13)),
], left=Inches(7.05), top=Inches(1.98), width=Inches(5.7), height=Inches(0.55))

code_block(s, [
    "P(n≥2 | n≥1) = 1 − μ·e^(−μ) / (1 − e^(−μ))",
    "",
    "   μ = 0.1   →    4.9 %",
    "   μ = 0.4   →   18.7 %",
    "   μ = 0.8   →   34.7 %",
    "   μ = 1.0   →   41.8 %   ← assumed in the previous report",
], Inches(7.05), Inches(2.55), Inches(5.7), Inches(2.05), size=11.5,
    accent=RED, bg=NOTEBG)

bullets(s, [
    ("Underwater the situation is worse than in air: channel loss is so large "
     "that a photon-number-splitting attacker can hide the extra loss they "
     "introduce. The controller clamps drive at MU_CAP = 12 of 15 even in "
     "CONSERVATIVE mode.", 0, dict(size=13.5)),
    ("A full defence needs decoy states — listed under future work.",
     0, dict(size=13.5, italic=True)),
], left=Inches(7.05), top=Inches(4.80), width=Inches(5.7),
    height=Inches(1.9), gap=8)

# ═══════════════════════════════════════════════════════════════════════════
section("IV", "Simulation Results",
        ("QBER", "P_click", "SKR", "Water types", "Wavelength", "Turbulence"))

# ── Status warnings ───────────────────────────────────────────────────────
s = new_slide("How these results were produced")
bullets(s, [
    ("Every figure in this part is a Monte-Carlo simulation of the SAME "
     "channel model that is compiled into the FPGA ROM "
     "(uwoc_channel_model.py → uwoc_lut_gen.py → uwoc_channel_rom.vh).",
     0, dict(size=15)),
    ("The simulation reproduces the RTL decision structure exactly: fading is "
     "drawn once per coherence block of 50,000 pulses, then click and error "
     "are decided by the same four comparisons the hardware performs.",
     0, dict(size=15)),
    ("Sample size is chosen per point to collect at least 3,000 sifted bits, "
     "so the Clopper–Pearson intervals shown are tight enough to read.",
     0, dict(size=15)),
], top=Inches(1.50), height=Inches(2.6), gap=12)

callout(s, "These are NOT hardware measurements. The FPGA board is currently "
           "faulty and the measurement campaign has not been run. What the "
           "figures establish is that the model, the fixed-point ROM "
           "generation and the decision logic are self-consistent and give "
           "physically sensible behaviour — the hardware campaign will test "
           "whether the built system matches them.",
        ML, Inches(4.30), MW, Inches(1.45), label="Please note:", size=14.5)

code_block(s, [
    "python python/report_sim_figs.py     regenerates every figure in Part IV",
    "python python/uwoc_channel_model.py  7 self-tests of the channel model",
], ML, Inches(6.02), MW, Inches(0.95), size=12.5)

# ── Results ───────────────────────────────────────────────────────────────
fig_slide("QBER vs distance", "fig04_qber_vs_distance.png",
          takeaways=[
              ("Water type sets the operating range: the 11 % Shor–Preskill "
               "bound is crossed at roughly 45 m in clear ocean, 15 m in "
               "coastal water and under 2 m in a turbid harbor.",
               0, dict(size=13.5)),
              ("The rise is driven by two effects at once — e_pol(d) grows "
               "with scattering optical depth, and as P_click falls toward "
               "Y₀ the background clicks (wrong half the time) take over, "
               "pushing QBER to 50 %.", 0, dict(size=13.5)),
          ])

fig_slide("Detection probability vs distance", "fig05_pclick_vs_distance.png",
          takeaways=[
              ("P_click falls from ~1e-2 to the noise floor Y₀ = 1.3e-5. Once "
               "the curve reaches that floor the link is delivering nothing "
               "but dark counts.", 0, dict(size=13.5)),
              ("This is the single most important number for the measurement "
               "plan: at 1e-3, collecting 1,000 sifted bits needs 2 million "
               "transmitted pulses.", 0, dict(size=13.5)),
          ])

fig_slide("Secure key rate vs distance", "fig06_skr_vs_distance.png",
          takeaways=[
              ("Asymptotic BB84 rate R = ½·f_rep·P_click·(1 − 2H₂(QBER)) at "
               "f_rep = 10 MHz. The cliff is where QBER reaches 11 % — beyond "
               "it no secure key exists at any sample size.",
               0, dict(size=13.5)),
              ("Clear ocean sustains tens of kbit/s out to ~25 m; the turbid "
               "harbor gives a usable key only within about a metre.",
               0, dict(size=13.5)),
          ])

fig_slide("Maximum secure distance across water types and turbulence",
          "fig10_max_secure_distance.png",
          takeaways=[
              ("Criterion: fewer than 10 % of monitoring windows exceed 11 % "
               "QBER — the outage criterion, which unlike the long-run mean "
               "does distinguish turbulence levels.", 0, dict(size=13.5)),
              ("A factor ~14 in extinction coefficient becomes a factor ~26 in "
               "secure distance. A turbid-harbor QKD link is a SHORT link: "
               "beyond a couple of metres the answer is to move the terminals, "
               "not to raise the drive level — that hits the PNS cap first.",
               0, dict(size=13.5)),
          ], side="top")

fig_slide("The optimal wavelength depends on turbidity",
          "fig09_wavelength.png",
          takeaways=[
              ("Clear ocean prefers blue (450 nm); as turbidity rises the "
               "optimum shifts toward red, and in the harbor 650 nm gives both "
               "the highest P_click and the lowest QBER.", 0, dict(size=13.5)),
              ("This is why λ is an adaptive knob rather than a design-time "
               "constant. Caveat: the λ-scaling of c(λ) is indicative, so the "
               "ordering is trustworthy but the exact crossover is not.",
               0, dict(size=13, color=INK2)),
          ], side="top")

fig_slide("Turbulence hides in the variance, not in the mean",
          "fig07_turbulence_mean_vs_std.png",
          takeaways=[
              ("Mean QBER moves from 3.00 % to 3.12 % across L1–L5 — "
               "essentially flat. The between-window standard deviation moves "
               "from 2.4 % to 14.5 %, a factor of six.", 0, dict(size=13.5)),
              ("Reported deliberately as a negative result: it is the "
               "experimental confirmation of Finding 2, and the justification "
               "for monitoring qber_jitter instead of mean QBER alone.",
               0, dict(size=13.5)),
          ], side="top")

fig_slide("Why the monitoring window must be large",
          "fig08_window_size.png",
          takeaways=[
              ("At 2⁸ and 2¹⁰ attempts, very weak and severe turbulence "
               "produce the SAME per-window QBER spread — the measurement is "
               "pure shot noise and carries no channel information.",
               0, dict(size=13.5)),
              ("At the deployed 2¹⁶ window the two separate cleanly (2.1 % vs "
               "16.9 %), which is exactly the signal the adaptive controller "
               "consumes as qber_jitter.", 0, dict(size=13.5)),
          ])

fig_slide("Sample size is the binding constraint",
          "fig12_sample_size.png",
          takeaways=[
              ("With 0 errors out of 93 sifted bits, the naive reading is "
               "\"QBER = 0 %, link secure\". The correct statement is "
               "QBER ≤ 3.2 % at 95 % confidence — and the model predicts "
               "1.64 % there, so zero errors happens 21 % of the time.",
               0, dict(size=13.5)),
              ("Methodology adopted: report a Clopper–Pearson interval on "
               "every QBER point; declare \"secure\" only when the 95 % UPPER "
               "BOUND stays below 11 %; plot points with n_sift < 16 as "
               "upper-limit arrows; budget n_sift ≥ 1000 per configuration.",
               0, dict(size=13.5)),
          ])

# ═══════════════════════════════════════════════════════════════════════════
section("V", "Verification, Status and Future Work",
        ("Verification plan", "Resources", "What is pending"))

# ── Verification ──────────────────────────────────────────────────────────
s = new_slide("Verification: RTL vs model vs measurement")
for i, (head, body, done) in enumerate([
    ("1 · RTL vs Python golden model — tb_uwoc_channel.v",
     "Click rate, error rate and frame integrity of the emulator checked "
     "against the analytical model in ModelSim.", True),
    ("2 · Closed-loop behaviour — tb_adaptive_loop.v",
     "Mode transitions, asymmetric hysteresis and wavelength hill-climbing. "
     "This testbench is what caught the 8-bit photon_rate saturation bug.",
     True),
    ("3 · Command path regression — tb_cmd_fifo.v",
     "Verifies that no qubit command is dropped when the PC streams commands "
     "back to back. An earlier version silently executed only 14 of every 32 "
     "commands, which depressed the measured P_click to 0.4× its true value.",
     True),
    ("4 · Hardware vs model — python/check_vs_theory.py",
     "Two statistical tests on measured data: clicks against "
     "Poisson(N·P_click_model) as a z-score, and errors against "
     "Binomial(n_sift, QBER_model) with a Clopper–Pearson interval.",
     False),
]):
    y = Inches(1.50 + i * 1.31)
    bg = RGBColor(0xF7, 0xF8, 0xFA) if done else NOTEBG
    accent = AQUA if done else RED
    _rect(s, ML, y, MW, Inches(1.14), bg)
    _rect(s, ML, y, Inches(0.055), Inches(1.14), accent)
    _, tf = _txbox(s, ML + Inches(0.28), y + Inches(0.13),
                   MW - Inches(1.9), Inches(0.95))
    _run(tf.paragraphs[0], head, size=14.5, bold=True, color=accent)
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    _run(p, body, size=12.5, color=INK)
    _, tf = _txbox(s, Inches(11.15), y + Inches(0.38), Inches(1.5), Inches(0.4))
    pp = tf.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    _run(pp, "PASSING" if done else "BLOCKED", size=12, bold=True, color=accent)

_, tf = _txbox(s, ML, Inches(6.62), MW, Inches(0.35))
_run(tf.paragraphs[0],
     "Layer 4 is blocked on the board fault. Layers 1–3 run in simulation and "
     "are unaffected.", size=12.5, color=RED, italic=True)

# ── Resources ─────────────────────────────────────────────────────────────
s = new_slide("Implementation results")
table(s, [
    ["Item", "Value"],
    ["Device", "Altera Cyclone II EP2C20F484C7"],
    ["Logic", "6,304 LE of 18,752   (34 %)"],
    ["Registers", "1,263"],
    ["Memory", "7,104 bits"],
    ["Clock", "50 MHz"],
    ["Toolchain", "Quartus II 13.0.1 — Analysis & Synthesis, Fitter and "
                  "Assembler all complete, 0 errors"],
], ML, Inches(1.55), Inches(7.4), col_w=[26, 74], size=13, head_size=12.5,
    row_h=Inches(0.37))

bullets(s, [
    ("Measurement throughput", 0, dict(bold=True, size=15.5)),
    ("86.6 qubit/s over UART at 115,200 bps, measured before the fault.",
     1, dict(size=13.5)),
    ("The limit is the PC-side round trip, not the FPGA: a qubit takes ~220 µs "
     "of hardware time (TX 50 µs + RX timeout 160 µs + gap 10 µs). Every "
     "non-detected qubit costs one host timeout — and underwater more than "
     "99 % of qubits are non-detected.", 1, dict(size=13.5)),
    ("Reporting every qubit, including the misses, removes this bottleneck "
     "and is the single highest-value next step for the measurement setup.",
     1, dict(size=13.5, color=RED)),
], left=Inches(8.15), top=Inches(1.55), width=Inches(4.6),
    height=Inches(4.6), gap=8)

callout(s, "Synthesis, fitting and timing are complete and clean. The fault "
           "is on the measurement path, not in the design as built.",
        ML, Inches(4.55), Inches(7.4), Inches(0.85), label="Status:", size=13.5)

# ── Future work ───────────────────────────────────────────────────────────
s = new_slide("Future work")
cols = [
    ("Measurement setup", BLUE, [
        "Report every qubit, including non-detections, to remove the "
        "host-timeout bottleneck (~5× throughput, enabling n_sift ≥ 1000 "
        "per configuration)",
        "Repair the board and run the full measurement campaign against the "
        "curves in Part IV",
        "Fit k_s in the polarization-error model against measured data "
        "instead of the current nominal value",
    ]),
    ("Protocol", AQUA, [
        "Decoy-state BB84 — the proper defence against PNS, which the "
        "drive-level cap only mitigates",
        "Finite-key security analysis rather than the asymptotic bound",
    ]),
    ("Channel model", GOLD, [
        "Bubbles and suspended particles as an explicit blockage process",
        "Measured a(λ), b(λ) spectra per water type",
    ]),
    ("System", RED, [
        "Two-FPGA Alice–Bob prototype over a real optical link",
        "One-Alice multi-Bob underwater sensor network, with a key "
        "management layer holding K(A,B1) … K(A,Bn)",
    ]),
]
for i, (head, color, items) in enumerate(cols):
    x = ML + Inches(i * 3.12)
    _rect(s, x, Inches(1.52), Inches(2.92), Inches(0.42), color)
    _, tf = _txbox(s, x + Inches(0.16), Inches(1.60), Inches(2.6), Inches(0.3))
    _run(tf.paragraphs[0], head, size=13.5, bold=True, color=WHITE)
    _rect(s, x, Inches(1.94), Inches(2.92), Inches(4.30),
          RGBColor(0xF8, 0xF8, 0xF6))
    _, tf = _txbox(s, x + Inches(0.16), Inches(2.10), Inches(2.6), Inches(4.0))
    for j, it in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        _run(p, "▪  ", size=12, color=color)
        _run(p, it, size=12, color=INK)

# ── Closing ───────────────────────────────────────────────────────────────
s = new_slide(None, footer=False)
_rect(s, Inches(0), Inches(0), SW, Inches(0.22), RED)
_rect(s, Inches(0), Inches(0.22), SW, Inches(0.06), GOLD)
_, tf = _txbox(s, ML, Inches(2.75), Inches(11.6), Inches(1.0))
_run(tf.paragraphs[0], "Thank you  —  questions welcome", size=38, bold=True,
     color=INK)
_rect(s, ML, Inches(3.95), Inches(2.1), Inches(0.045), RED)
_, tf = _txbox(s, ML, Inches(4.30), Inches(11.6), Inches(1.2))
_run(tf.paragraphs[0], "Lê Công Khánh   ·   Đào Đặng Minh Hoàng",
     size=17, color=INK2)
p = tf.add_paragraph()
p.space_before = Pt(10)
_run(p, "Reproduce every number in this deck:", size=13, color=INK3)
p = tf.add_paragraph()
_run(p, "python python/uwoc_channel_model.py   ·   "
        "python python/report_sim_figs.py", size=13, color=BLUE, font=MONO)
_rect(s, Inches(0), Inches(7.28), SW, Inches(0.22), RED)

# ── Appendix ──────────────────────────────────────────────────────────────
s = new_slide("Appendix — constants and where they come from")
rows_l = [
    ["Quantity", "Value", "Source"],
    ["μ", "0.1 photon/pulse", "uwoc_channel_model.py:169"],
    ["η_det", "0.18", ":159"],
    ["dark / background", "60 Hz / 200 Hz", ":160-161"],
    ["gate", "50 ns", ":162"],
    ["Y₀", "1.30e-5", ":181-183"],
    ["f_rep", "10 MHz", ":163"],
    ["τ_coh", "5 ms → 50,000 pulses", ":165, :186-188"],
    ["D_rx", "50.8 mm", ":153"],
    ["θ_div", "1e-3 rad", ":154"],
    ["F", "0.85", ":158"],
    ["e₀ / k_s", "0.01 / 0.04", ":170-171"],
]
rows_r = [
    ["Quantity", "Value", "Source"],
    ["c: clear/coastal/harbor", "0.151 / 0.398 / 2.190 m⁻¹", ":92-103"],
    ["λ", "450 / 532 / 650 nm", ":117"],
    ["σ²_ho L1–L5", "0.02 / 0.08 / 0.30 / 1.00 / 3.00", ":119-138"],
    ["COH_LOG2", "18 → 5.24 ms", "uwoc_channel.v:53"],
    ["h scale", "12 bit, 256 = 1.0, max 4095", "uwoc_channel.v:92-93"],
    ["probabilities", "24 bit, scale 2²⁴", "uwoc_channel.v:94"],
    ["p_noise", "218 / 2²⁴", "top_module.v:291"],
    ["monitoring window", "2¹⁶ attempts", "channel_monitor.v:47"],
    ["MIN_SIFT", "16", "channel_monitor.v:58"],
    ["QBER thresholds", "4 / 8 / 11 / 15 %", "adaptive_controller.v:91-94"],
    ["SNR / jitter thresholds", "160 / 96 / 40  ·  6 / 16", ":97-99, :113-114"],
    ["MU_CAP", "12", ":45"],
]
table(s, rows_l, ML, Inches(1.50), Inches(5.9), col_w=[33, 34, 33],
      size=10.5, head_size=10.5, row_h=Inches(0.315))
table(s, rows_r, Inches(6.85), Inches(1.50), Inches(5.9), col_w=[33, 36, 31],
      size=10.5, head_size=10.5, row_h=Inches(0.315))

# ═══════════════════════════════════════════════════════════════════════════
prs.save(OUTFILE)
print(f"saved  {OUTFILE}")
print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
